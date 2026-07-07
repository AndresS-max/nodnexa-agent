"""Extracción de contenido por formato (etapa 2 del pipeline RAG).

Cada loader devuelve una lista de Documents de LangChain con el texto limpio
y metadatos: archivo de origen, categoría de negocio, formato y ubicación
(página/fila/sección) para poder citar fuentes en las respuestas.
"""
import json
import re
from pathlib import Path

import pandas as pd
from langchain_core.documents import Document
from pypdf import PdfReader

from src.config import CATEGORIA_POR_ARCHIVO


def _limpiar(texto: str) -> str:
    """Elimina ruido común de extracción: espacios duplicados y líneas vacías."""
    texto = re.sub(r"[ \t]+", " ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto.strip()


def _meta(path: Path, formato: str, **extra) -> dict:
    return {
        "archivo": path.name,
        "categoria": CATEGORIA_POR_ARCHIVO.get(path.name, "General"),
        "formato": formato,
        **extra,
    }


def load_pdf(path: Path) -> list[Document]:
    """Un Document por página, para citar 'archivo, página N'."""
    docs = []
    reader = PdfReader(str(path))
    for num, page in enumerate(reader.pages, start=1):
        texto = _limpiar(page.extract_text() or "")
        if texto:
            docs.append(Document(page_content=texto,
                                 metadata=_meta(path, "pdf", pagina=num)))
    return docs


def load_csv(path: Path) -> list[Document]:
    """Cada fila se convierte en una frase legible con sus encabezados,
    porque las tablas pierden el sentido si se tratan como texto corrido."""
    df = pd.read_csv(path)
    docs = []
    for idx, row in df.iterrows():
        partes = [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
        texto = ". ".join(str(p) for p in partes)
        docs.append(Document(page_content=texto,
                             metadata=_meta(path, "csv", fila=int(idx) + 2)))
    return docs


def load_markdown(path: Path) -> list[Document]:
    """Divide por secciones (##) para preservar la estructura lógica."""
    contenido = path.read_text(encoding="utf-8")
    secciones = re.split(r"\n(?=## )", contenido)
    docs = []
    for seccion in secciones:
        texto = _limpiar(seccion)
        if not texto:
            continue
        titulo = texto.splitlines()[0].lstrip("# ").strip()
        docs.append(Document(page_content=texto,
                             metadata=_meta(path, "markdown", seccion=titulo)))
    return docs


def load_json(path: Path) -> list[Document]:
    """Convierte cada caso de uso en un párrafo comprensible."""
    data = json.loads(path.read_text(encoding="utf-8"))
    docs = []
    for caso in data.get("casos_de_uso", []):
        texto = (
            f"Industria: {caso['industria']}.\n"
            f"Problemas típicos: {'; '.join(caso['problemas_tipicos'])}.\n"
            f"Solución de Nodnexa: {caso['solucion_nodnexa']}.\n"
            f"Servicios relacionados: {', '.join(caso['servicios_relacionados'])}.\n"
            f"Resultado esperado: {caso['resultado_esperado']}."
        )
        docs.append(Document(page_content=texto,
                             metadata=_meta(path, "json", seccion=caso["industria"])))
    return docs


def load_docx(path: Path) -> list[Document]:
    """Word: párrafos agrupados bajo su encabezado más cercano, y las tablas
    convertidas fila→frase (misma lógica que el CSV)."""
    from docx import Document as DocxDocument

    docx = DocxDocument(str(path))
    docs, seccion, buffer = [], "Inicio", []

    def cerrar_seccion():
        texto = _limpiar("\n".join(buffer))
        if texto:
            docs.append(Document(page_content=f"{seccion}\n{texto}",
                                 metadata=_meta(path, "docx", seccion=seccion)))

    for p in docx.paragraphs:
        if p.style.name.startswith("Heading") and p.text.strip():
            cerrar_seccion()
            seccion, buffer = p.text.strip(), []
        elif p.text.strip():
            buffer.append(p.text.strip())
    cerrar_seccion()

    for num, tabla in enumerate(docx.tables, start=1):
        filas = [[c.text.strip() for c in fila.cells] for fila in tabla.rows]
        if len(filas) < 2:
            continue
        encabezados = filas[0]
        for fila in filas[1:]:
            texto = ". ".join(f"{h}: {v}" for h, v in zip(encabezados, fila) if v)
            if texto:
                docs.append(Document(page_content=texto,
                                     metadata=_meta(path, "docx",
                                                    seccion=f"tabla {num}")))
    return docs


def load_xlsx(path: Path) -> list[Document]:
    """Excel: cada fila de cada hoja se convierte en frase con encabezados
    (las planillas tienen lógica tabular, no de texto corrido)."""
    hojas = pd.read_excel(path, sheet_name=None)
    docs = []
    for nombre_hoja, df in hojas.items():
        for idx, row in df.iterrows():
            partes = [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
            if partes:
                docs.append(Document(
                    page_content=". ".join(str(p) for p in partes),
                    metadata=_meta(path, "xlsx", seccion=nombre_hoja,
                                   fila=int(idx) + 2),
                ))
    return docs


def load_html(path: Path) -> list[Document]:
    """HTML: texto sin etiquetas ni scripts, dividido por encabezados h1-h3."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"),
                         "html.parser")
    for basura in soup(["script", "style", "nav", "footer"]):
        basura.decompose()

    docs, seccion, buffer = [], soup.title.string.strip() if soup.title and soup.title.string else "Inicio", []

    def cerrar_seccion():
        texto = _limpiar(" ".join(buffer))
        if texto:
            docs.append(Document(page_content=f"{seccion}\n{texto}",
                                 metadata=_meta(path, "html", seccion=seccion)))

    cuerpo = soup.body or soup
    for elem in cuerpo.find_all(["h1", "h2", "h3", "p", "li", "td", "th"]):
        texto = elem.get_text(" ", strip=True)
        if not texto:
            continue
        if elem.name in ("h1", "h2", "h3"):
            cerrar_seccion()
            seccion, buffer = texto, []
        else:
            buffer.append(texto)
    cerrar_seccion()
    return docs


def load_pptx(path: Path) -> list[Document]:
    """PowerPoint: un Document por diapositiva, incluyendo las notas del
    orador (suelen contener contexto importante)."""
    from pptx import Presentation

    pres = Presentation(str(path))
    docs = []
    for num, slide in enumerate(pres.slides, start=1):
        textos = [shape.text.strip() for shape in slide.shapes
                  if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if slide.has_notes_slide:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                textos.append(f"Notas del orador: {notas}")
        texto = _limpiar("\n".join(textos))
        if texto:
            docs.append(Document(page_content=texto,
                                 metadata=_meta(path, "pptx", diapositiva=num)))
    return docs


LOADERS = {
    ".pdf": load_pdf,
    ".csv": load_csv,
    ".md": load_markdown,
    ".json": load_json,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".html": load_html,
    ".htm": load_html,
    ".pptx": load_pptx,
}


def load_documents(docs_dir: Path) -> list[Document]:
    """Carga todos los documentos soportados de la carpeta."""
    documentos = []
    for path in sorted(docs_dir.iterdir()):
        loader = LOADERS.get(path.suffix.lower())
        if loader:
            documentos.extend(loader(path))
    return documentos
