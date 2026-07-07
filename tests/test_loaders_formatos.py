"""Pruebas de los loaders de Word, Excel, HTML y PowerPoint.

Cada test genera un archivo real del formato en un directorio temporal
y verifica extracción, metadatos y ubicación citable.
"""
import pandas as pd
import pytest

from src.ingestion.loaders import load_docx, load_html, load_pptx, load_xlsx


@pytest.fixture
def tmp_docs(tmp_path):
    return tmp_path


def test_docx_extrae_secciones_y_tablas(tmp_docs):
    from docx import Document as DocxDocument

    f = tmp_docs / "politica.docx"
    d = DocxDocument()
    d.add_heading("Política de viáticos", level=1)
    d.add_paragraph("Los viáticos se reembolsan en un máximo de 15 días.")
    tabla = d.add_table(rows=2, cols=2)
    tabla.rows[0].cells[0].text = "Concepto"
    tabla.rows[0].cells[1].text = "Tope"
    tabla.rows[1].cells[0].text = "Hotel"
    tabla.rows[1].cells[1].text = "USD 80"
    d.save(str(f))

    docs = load_docx(f)
    assert any("viáticos se reembolsan" in x.page_content for x in docs)
    assert any(x.metadata["seccion"] == "Política de viáticos" for x in docs)
    assert any("Concepto: Hotel. Tope: USD 80" in x.page_content for x in docs)
    assert all(x.metadata["formato"] == "docx" for x in docs)


def test_xlsx_convierte_filas_por_hoja(tmp_docs):
    f = tmp_docs / "inventario.xlsx"
    pd.DataFrame({"SKU": ["A1", "A2"], "Stock": [10, 5]}).to_excel(
        f, sheet_name="Bodega", index=False)

    docs = load_xlsx(f)
    assert len(docs) == 2
    assert "SKU: A1. Stock: 10" in docs[0].page_content
    assert docs[0].metadata["seccion"] == "Bodega"
    assert docs[0].metadata["fila"] == 2  # fila 2 del Excel (1 es encabezado)


def test_html_limpia_y_divide_por_encabezados(tmp_docs):
    f = tmp_docs / "faq.html"
    f.write_text(
        "<html><head><title>FAQ</title><script>var x=1;</script></head>"
        "<body><h2>Envíos</h2><p>Entregamos en 3 días.</p>"
        "<h2>Pagos</h2><p>Aceptamos tarjeta.</p></body></html>",
        encoding="utf-8")

    docs = load_html(f)
    assert len(docs) == 2
    assert "Entregamos en 3 días" in docs[0].page_content
    assert docs[1].metadata["seccion"] == "Pagos"
    assert all("var x" not in x.page_content for x in docs)  # sin scripts


def test_pptx_incluye_notas_del_orador(tmp_docs):
    from pptx import Presentation
    from pptx.util import Inches

    f = tmp_docs / "pitch.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    slide.shapes.title.text = "Plan comercial"
    slide.notes_slide.notes_text_frame.text = "Enfatizar el ROI del cliente."
    pres.save(str(f))

    docs = load_pptx(f)
    assert len(docs) == 1
    assert "Plan comercial" in docs[0].page_content
    assert "Notas del orador: Enfatizar el ROI" in docs[0].page_content
    assert docs[0].metadata["diapositiva"] == 1
